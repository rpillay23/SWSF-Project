import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Your investment data
data = [
    # Traditional
    {"Investment Name": "Large Cap Equities", "Category": "Equities", "Expected Return (%)": 8, "Risk Level (1-10)": 6, "Cap Rate (%)": None, "Liquidity": "High", "Min Investment": 0, "Fund Manager": "Vanguard"},
    {"Investment Name": "Government Bonds", "Category": "Fixed Income", "Expected Return (%)": 4, "Risk Level (1-10)": 2, "Cap Rate (%)": None, "Liquidity": "High", "Min Investment": 1000, "Fund Manager": "BlackRock"},
    {"Investment Name": "Money Market Funds", "Category": "Cash and Equivalents", "Expected Return (%)": 2, "Risk Level (1-10)": 1, "Cap Rate (%)": None, "Liquidity": "High", "Min Investment": 100, "Fund Manager": "Fidelity"},
    {"Investment Name": "Forex Managed Account", "Category": "Currencies (Forex)", "Expected Return (%)": 6, "Risk Level (1-10)": 7, "Cap Rate (%)": None, "Liquidity": "High", "Min Investment": 5000, "Fund Manager": "OANDA"},

    # Alternatives
    {"Investment Name": "REIT - Commercial", "Category": "Real Estate", "Expected Return (%)": 9, "Risk Level (1-10)": 5, "Cap Rate (%)": 5.7, "Liquidity": "Medium", "Min Investment": 25000, "Fund Manager": "Prologis"},
    {"Investment Name": "Gold ETF", "Category": "Commodities", "Expected Return (%)": 6, "Risk Level (1-10)": 4, "Cap Rate (%)": None, "Liquidity": "High", "Min Investment": 1000, "Fund Manager": "SPDR"},
    {"Investment Name": "Life Settlement Fund", "Category": "Life Settlements", "Expected Return (%)": 10, "Risk Level (1-10)": 6, "Cap Rate (%)": None, "Liquidity": "Low", "Min Investment": 100000, "Fund Manager": "Abacus Life"},
    {"Investment Name": "Mezzanine Finance Fund", "Category": "Mezzanine Financing", "Expected Return (%)": 12, "Risk Level (1-10)": 7, "Cap Rate (%)": None, "Liquidity": "Low", "Min Investment": 50000, "Fund Manager": "Golub Capital"},
    {"Investment Name": "Direct Lending Fund", "Category": "Direct Lending", "Expected Return (%)": 11, "Risk Level (1-10)": 6, "Cap Rate (%)": None, "Liquidity": "Low", "Min Investment": 50000, "Fund Manager": "Ares Capital"},
    {"Investment Name": "Raw Land Development", "Category": "Raw Developable Land", "Expected Return (%)": 14, "Risk Level (1-10)": 8, "Cap Rate (%)": 7.5, "Liquidity": "Low", "Min Investment": 100000, "Fund Manager": "Private Group"},
    {"Investment Name": "Fund of Funds - Alt Blend", "Category": "Fund of Funds", "Expected Return (%)": 9, "Risk Level (1-10)": 5, "Cap Rate (%)": None, "Liquidity": "Medium", "Min Investment": 25000, "Fund Manager": "Partners Group"},
    {"Investment Name": "Infrastructure Fund", "Category": "Infrastructure", "Expected Return (%)": 8.5, "Risk Level (1-10)": 4, "Cap Rate (%)": 6.0, "Liquidity": "Medium", "Min Investment": 50000, "Fund Manager": "Macquarie"}
]  # Use your full investment_data list here from the previous script

df = pd.DataFrame(data)

# === PART 1: SAVE EXCEL FILE ===
df.to_excel("Comprehensive_Investment_Matrix.xlsx", index=False)

# === PART 2: CREATE CHART ===
plt.figure(figsize=(10, 5))
plt.bar(df["Investment Name"], df["Expected Return (%)"], color="steelblue")
plt.xticks(rotation=90)
plt.ylabel("Expected Return (%)")
plt.title("Expected Return by Investment")
plt.tight_layout()
chart_path = "expected_returns_chart.png"
plt.savefig(chart_path)
plt.close()

# === PART 3: CREATE POWERPOINT ===
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Comprehensive Investment Overview"
slide.placeholders[1].text = "HNW Client Alternative & Traditional Investment Matrix"

# Averages slide
avg = df[["Expected Return (%)", "Risk Level (1-10)", "Cap Rate (%)"]].mean(numeric_only=True).round(2)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Portfolio Averages"
content = slide.placeholders[1]
content.text = (
    f"Expected Return (avg): {avg['Expected Return (%)']}%\n"
    f"Risk Level (avg): {avg['Risk Level (1-10)']}\n"
    f"Cap Rate (avg): {avg['Cap Rate (%)']}"
)

# Chart slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Expected Return by Investment"
slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))

# Investment category list
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Investment Categories Included"
text_frame = slide.placeholders[0].text_frame
for cat in df["Category"].unique():
    p = text_frame.add_paragraph()
    p.text = f"- {cat}"

# Save PPT
prs.save("HNW_Investment_Presentation.pptx")
print("✅ PowerPoint created: HNW_Investment_Presentation.pptx")