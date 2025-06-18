import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import yfinance as yf

# --- PAGE CONFIG ---
st.set_page_config(page_title="Automated Investment Matrix", layout="wide")

# --- CUSTOM CSS FOR THEME AND LAYOUT ---
st.markdown("""
    <style>
    /* Reset and fonts */
    html, body, [class*="css"] {
        font-family: 'Helvetica Neue', Helvetica, sans-serif;
        background-color: white;
        color: #111111;
    }
    /* HEADER */
    .app-header {
        background-color: #111111;
        padding: 10px 20px;
        margin: 0;
        width: 100vw;
        position: relative;
        left: calc(-50vw + 50%);
        top: 0;
        user-select: none;
        border-radius: 0;
    }
    .app-header h1 {
        color: white;
        font-size: 22px;
        margin: 0;
        font-weight: 700;
        line-height: 1.2;
    }
    .app-header p {
        color: #f44336;
        font-size: 12px;
        margin: 4px 0 0 0;
        font-weight: 500;
    }
    /* LEFT PERMANENT SIDEBAR */
    #market-indices {
        position: fixed;
        top: 80px;  /* just below header */
        left: 0;
        width: 260px;
        height: calc(100vh - 80px);
        background-color: #111111;
        color: white;
        overflow-y: auto;
        padding: 15px 10px;
        z-index: 100;
        border-right: 3px solid #f44336;
        font-size: 13px;
    }
    #market-indices h2 {
        color: #f44336;
        font-weight: 700;
        margin-bottom: 10px;
        font-size: 18px;
    }
    #market-indices .index-name {
        font-weight: 600;
        margin-top: 10px;
        color: white;
    }
    #market-indices .metric {
        margin-left: 5px;
        color: #f44336;
        font-weight: 700;
    }
    #market-indices .delta-positive {
        color: #4caf50;
        font-weight: 700;
    }
    #market-indices .delta-negative {
        color: #f44336;
        font-weight: 700;
    }
    #market-indices canvas {
        margin-top: 5px;
        margin-bottom: 15px;
    }
    /* MAIN CONTENT */
    #main-content {
        margin-left: 280px;
        padding: 15px 25px 50px 25px;
        max-width: calc(100vw - 280px);
    }
    /* Smaller charts and fonts */
    .small-chart {
        width: 260px !important;
        height: 160px !important;
    }
    .small-metric {
        font-size: 12px !important;
        padding: 0.3rem 0.5rem !important;
    }
    /* Buttons */
    .stButton > button {
        background-color: #111111;
        color: #f44336;
        border-radius: 6px;
        padding: 0.4em 1em;
        border: 2px solid #f44336;
        font-weight: 700;
        transition: background-color 0.3s ease, color 0.3s ease;
        font-size: 14px;
    }
    .stButton > button:hover {
        background-color: #f44336;
        color: #111111;
        border: 2px solid #f44336;
    }
    /* Data editor */
    .dataframe-container {
        max-height: 250px;
        overflow-y: auto;
    }
    </style>

    <div class="app-header">
        <h1>Automated Investment Matrix</h1>
        <p>Modular Investment Analysis Platform with Real Market Data for Portfolio Optimization and Financial Advisory.<br>
           Designed for portfolio managers and finance professionals to generate, analyze, and optimize investment portfolios using real-time data.</p>
    </div>
""", unsafe_allow_html=True)

# --- FUNCTIONS ---

@st.cache_data(ttl=600)
def get_index_data(ticker):
    try:
        index = yf.Ticker(ticker)
        hist = index.history(period="1mo")
        hist.reset_index(inplace=True)
        return hist
    except Exception:
        return pd.DataFrame()

def format_market_metric(latest, prev):
    latest_close = latest['Close']
    prev_close = prev['Close']
    change = latest_close - prev_close
    pct_change = (change / prev_close) * 100
    delta_class = "delta-positive" if change >= 0 else "delta-negative"
    return latest_close, change, pct_change, delta_class

def plot_mini_line_chart(data):
    fig, ax = plt.subplots(figsize=(3, 1))
    ax.plot(data['Date'], data['Close'], color='#f44336', linewidth=1.7)
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    plt.tight_layout()
    st.pyplot(fig, use_container_width=True)

def sanitize_string(s):
    if isinstance(s, str):
        return (
            s.replace("–", "-")
             .replace("’", "'")
             .replace("“", '"')
             .replace("”", '"')
             .replace("•", "-")
             .replace("©", "(c)")
        )
    return s

# --- LOAD DATA ---

try:
    df = pd.read_excel("Comprehensive_Investment_Matrix.xlsx")
    df = df.applymap(sanitize_string)
except Exception as e:
    st.error(f"Error loading data file: {e}")
    st.stop()

# --- LEFT PERMANENT MARKET INDICES ---

st.markdown('<div id="market-indices">', unsafe_allow_html=True)
st.markdown("<h2>Real-Time Market Indices</h2>")

for ticker, name in [("^GSPC", "S&P 500"), ("^IXIC", "Nasdaq"), ("^DJI", "Dow Jones")]:
    data = get_index_data(ticker)
    if data.empty:
        st.markdown(f"<div class='index-name'>{name}</div><div>Data unavailable</div>", unsafe_allow_html=True)
        continue

    latest = data.iloc[-1]
    prev = data.iloc[-2]
    latest_close, change, pct_change, delta_class = format_market_metric(latest, prev)

    st.markdown(f"<div class='index-name'>{name}</div>", unsafe_allow_html=True)
    st.markdown(f"<span class='metric'>${latest_close:,.2f}</span> <span class='{delta_class}'>{change:+.2f} ({pct_change:+.2f}%)</span>", unsafe_allow_html=True)
    plot_mini_line_chart(data)

st.markdown('</div>', unsafe_allow_html=True)

# --- MAIN CONTENT ---

with st.container():
    st.markdown('<div id="main-content">', unsafe_allow_html=True)

    # Investment Generator filters

    st.header("Investment Generator")

    # Filter by Category (Inflation Hedge, etc)
    categories = sorted(df['Category'].dropna().unique())
    selected_categories = st.multiselect("Select Investment Categories:", categories, default=categories)

    # Slider for Minimum Investment
    min_inv_default = int(df['Minimum Investment ($)'].min())
    max_inv_default = int(df['Minimum Investment ($)'].max())
    min_inv = st.slider("Minimum Investment ($)", min_value=min_inv_default, max_value=max_inv_default, value=min_inv_default, step=1000)

    # Slider for Expected Return %
    exp_ret_min = float(df['Expected Return (%)'].min())
    exp_ret_max = float(df['Expected Return (%)'].max())
    exp_ret = st.slider("Minimum Expected Return (%)", min_value=exp_ret_min, max_value=exp_ret_max, value=exp_ret_min, step=0.1)

    # Slider for Risk Level (1-10)
    risk_min = int(df['Risk Level (1-10)'].min())
    risk_max = int(df['Risk Level (1-10)'].max())
    max_risk = st.slider("Maximum Risk Level (1-10)", min_value=risk_min, max_value=risk_max, value=risk_max, step=1)

    # Filter dataframe accordingly
    filtered_df = df[
        (df['Category'].isin(selected_categories)) &
        (df['Minimum Investment ($)'] >= min_inv) &
        (df['Expected Return (%)'] >= exp_ret) &
        (df['Risk Level (1-10)'] <= max_risk)
    ]

    st.write(f"### {len(filtered_df)} Investments Matching Your Criteria")
    st.dataframe(filtered_df.reset_index(drop=True), height=250)

    st.divider()

    # Portfolio overview metrics
    st.header("Portfolio Summary Metrics")
    c1, c2, c3, c4, c5, c6, c7 = st.columns(7)
    c1.metric("Avg Return (%)", f"{filtered_df['Expected Return (%)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c2.metric("Avg Risk (1–10)", f"{filtered_df['Risk Level (1-10)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c3.metric("Avg Cap Rate (%)", f"{filtered_df['Cap Rate (%)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c4.metric("Avg Liquidity", f"{filtered_df['Liquidity (1–10)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c5.metric("Avg Volatility", f"{filtered_df['Volatility (1–10)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c6.metric("Avg Fees (%)", f"{filtered_df['Fees (%)'].mean():.2f}" if not filtered_df.empty else "N/A")
    c7.metric("Avg Min Investment", f"${filtered_df['Minimum Investment ($)'].mean():,.0f}" if not filtered_df.empty else "N/A")

    st.divider()

    # Visual Insights: 4 charts side-by-side smaller

    st.header("Visual Insights")
    ch1, ch2, ch3, ch4 = st.columns(4)

    def plot_bar(ax, x, y, xlabel, ylabel):
        ax.bar(x, y, color="#f44336")
        ax.tick_params(axis='x', rotation=45, labelsize=7)
        ax.tick_params(axis='y', labelsize=7)
        ax.set_xlabel(xlabel, fontsize=8)
        ax.set_ylabel(ylabel, fontsize=8)

    with ch1:
        st.markdown("**Expected Return (%)**")
        fig, ax = plt.subplots(figsize=(3, 2))
        plot_bar(ax, filtered_df["Investment Name"], filtered_df["Expected Return (%)"], "", "Return (%)")
        plt.tight_layout()
        st.pyplot(fig)

    with ch2:
        st.markdown("**Liquidity vs Volatility**")
        fig, ax = plt.subplots(figsize=(3, 2))
        ax.scatter(filtered_df["Volatility (1–10)"], filtered_df["Liquidity (1–10)"], s=30,
                   c='red', alpha=0.7, edgecolors='black')
        ax.set_xlabel("Volatility", fontsize=8)
        ax.set_ylabel("Liquidity", fontsize=8)
        ax.tick_params(axis='x', labelsize=7)
        ax.tick_params(axis='y', labelsize=7)
        plt.tight_layout()
        st.pyplot(fig)

    with ch3:
        st.markdown("**Fees (%) vs Expected Return (%)**")
        fig, ax = plt.subplots(figsize=(3, 2))
        ax.scatter(filtered_df["Fees (%)"], filtered_df["Expected Return (%)"], s=30,
                   c='red', alpha=0.7, edgecolors='black')
        ax.set_xlabel("Fees (%)", fontsize=8)
        ax.set_ylabel("Return (%)", fontsize=8)
        ax.tick_params(axis='x', labelsize=7)
        ax.tick_params(axis='y', labelsize=7)
        plt.tight_layout()
        st.pyplot(fig)

    with ch4:
        st.markdown("**Risk Level Distribution**")
        fig, ax = plt.subplots(figsize=(3, 2))
        ax.hist(filtered_df["Risk Level (1-10)"], bins=10, color="#f44336", alpha=0.8)
        ax.set_xlabel("Risk Level", fontsize=8)
        ax.set_ylabel("Count", fontsize=8)
        ax.tick_params(axis='x', labelsize=7)
        ax.tick_params(axis='y', labelsize=7)
        plt.tight_layout()
        st.pyplot(fig)

    st.divider()

    # Export report options

    st.header("Export Reports")

    export_format = st.radio("Select export format", ["PowerPoint", "Word Document"], horizontal=True)

    def generate_ppt_report(df_report):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Investment Portfolio Report"

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = f"Average Expected Return: {df_report['Expected Return (%)'].mean():.2f}%"
        tf.add_paragraph().text = f"Average Risk Level: {df_report['Risk Level (1-10)'].mean():.2f}"
        tf.add_paragraph().text = f"Average Fees: {df_report['Fees (%)'].mean():.2f}%"

        filename = "portfolio_report.pptx"
        prs.save(filename)
        return filename

    def generate_word_report(df_report):
        from docx import Document

        doc = Document()
        doc.add_heading("Investment Portfolio Report", 0)
        doc.add_paragraph(f"Average Expected Return: {df_report['Expected Return (%)'].mean():.2f}%")
        doc.add_paragraph(f"Average Risk Level: {df_report['Risk Level (1-10)'].mean():.2f}")
        doc.add_paragraph(f"Average Fees: {df_report['Fees (%)'].mean():.2f}%")
        filename = "portfolio_report.docx"
        doc.save(filename)
        return filename

    if st.button("Generate Report"):
        if filtered_df.empty:
            st.warning("No investments selected to generate report.")
        else:
            if export_format == "PowerPoint":
                ppt_file = generate_ppt_report(filtered_df)
                with open(ppt_file, "rb") as f:
                    st.download_button("Download PowerPoint Report", f, ppt_file)
            else:
                doc_file = generate_word_report(filtered_df)
                with open(doc_file, "rb") as f:
                    st.download_button("Download Word Report", f, doc_file)

    st.markdown('</div>', unsafe_allow_html=True)
